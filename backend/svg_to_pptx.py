"""
SVG to PPTX converter - converts vector graphics from SVG to python-pptx shapes.
Handles: rect, circle, ellipse, line, polygon, path (basic)
"""
from xml.etree import ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re, math

# SVG namespace
NS = 'http://www.w3.org/2000/svg'

# 1280×720 viewBox → 13.333×7.5 inches
SVG_W = 1280; SVG_H = 720
PPT_W = Inches(13.333); PPT_H = Inches(7.5)
SX = 13.333 / SVG_W; SY = 7.5 / SVG_H

def svg_to_pptx(svg_path, pptx_path=None):
    prs = Presentation()
    prs.slide_width = PPT_W; prs.slide_height = PPT_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tree = ET.parse(svg_path)
    root = tree.getroot()

    # Parse SVG dimensions
    vb = root.get('viewBox', '0 0 1280 720').split()
    w, h = int(vb[2]), int(vb[3])

    # Build gradient/fill lookup from <defs>
    defs = {}
    for defn in root.findall(f'{{{NS}}}defs') or []:
        for el in defn:
            tag = el.tag.replace(f'{{{NS}}}', '')
            id_ = el.get('id')
            if id_:
                defs[id_] = el

    # Process all graphic elements
    for el in root:
        tag = el.tag.replace(f'{{{NS}}}', '')
        if tag in ('rect', 'circle', 'ellipse', 'line', 'polygon', 'polyline', 'path'):
            process_element(slide, el, tag, defs, w, h)

    if pptx_path is None:
        import os
        base = os.path.splitext(os.path.basename(svg_path))[0]
        pptx_path = base + '.pptx'
    prs.save(pptx_path)
    return pptx_path


def px(val, ref_w):
    """Convert SVG length to EMU"""
    if isinstance(val, str):
        val = val.strip()
        if val.endswith('%'):
            return None  # percentage - skip for now
        if val.endswith('px'):
            v = float(val[:-2])
        elif val.endswith('pt'):
            v = float(val[:-2])
        else:
            try:
                v = float(val)
            except:
                return None
    else:
        v = float(val)
    return int(v * 914400 / ref_w)


def get_fill(el, defs):
    """Resolve fill color from SVG element"""
    fill = el.get('fill', '#000000')
    if fill == 'none':
        return None
    if fill.startswith('url(#'):
        grad_id = fill[5:-1]
        grad = defs.get(grad_id)
        if grad is not None:
            # Return gradient element for special handling
            return grad
    return parse_color(fill)


def parse_color(c):
    """Parse SVG color to RGBColor"""
    if c is None or c == 'none' or c == 'transparent':
        return None
    c = c.strip()
    if c.startswith('#'):
        c = c[1:]
        if len(c) == 3:
            c = c[0]*2 + c[1]*2 + c[2]*2
        if len(c) == 6:
            return RGBColor(int(c[0:2],16), int(c[2:4],16), int(c[4:6],16))
    return None


def get_opacity(el):
    o = el.get('opacity', '1')
    try:
        return float(o)
    except:
        return 1.0


def add_pptx_shape(slide, etype, x, y, w, h, fill_rgb, stroke_rgb=None, stroke_w=0):
    """Add a shape to the slide"""
    # python-pptx shape types: 1=RECTANGLE, 9=OVAL, etc.
    shape = slide.shapes.add_shape(etype, x, y, w, h)
    if fill_rgb:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if stroke_rgb:
        shape.line.color.rgb = stroke_rgb
        shape.line.width = Emu(stroke_w)
    else:
        shape.line.fill.background()
    return shape


def process_element(slide, el, tag, defs, svg_w, svg_h):
    """Convert SVG element to PPTX shape"""

    # Scale factor from SVG px to EMU
    # EMU per inch = 914400, slide width in inches = 13.333
    def to_emu(v, ref): return int(v * 914400 * 13.333 / ref)

    opacity = get_opacity(el)
    fill = el.get('fill', None)
    stroke = el.get('stroke', 'none')
    sw_str = el.get('stroke-width', '0')
    sw = 0
    try:
        sw = float(sw_str)
    except:
        pass

    sfill = None
    sstroke = None
    if fill and fill != 'none':
        sfill = parse_color(fill)
    if stroke and stroke != 'none':
        sstroke = parse_color(stroke)

    # Adjust opacity for fill
    if opacity < 1.0 and sfill:
        sfill = adjust_opacity(sfill, opacity)

    x = to_emu(0, svg_w); y = to_emu(0, svg_h)
    w = to_emu(1280, svg_w); h = to_emu(720, svg_h)

    if tag == 'rect':
        x = to_emu(float(el.get('x', 0)), svg_w)
        y = to_emu(float(el.get('y', 0)), svg_h)
        w = to_emu(float(el.get('width', 0)), svg_w)
        h = to_emu(float(el.get('height', 0)), svg_h)
        rx_str = el.get('rx', '0'); ry_str = el.get('ry', '0')
        try: rx = float(rx_str)
        except: rx = 0
        try: ry = float(ry_str)
        except: ry = 0

        # Use RECTANGLE type - python-pptx doesn't natively support rounded rects
        # but we use MSO_SHAPE_TYPE.RECTANGLE for now
        shape = slide.shapes.add_shape(1, x, y, w, h)
        if sfill:
            shape.fill.solid(); shape.fill.fore_color.rgb = sfill
        else:
            shape.fill.background()
        if sstroke:
            shape.line.color.rgb = sstroke
            shape.line.width = Emu(int(sw * 914400 * 13.333 / svg_w))
        else:
            shape.line.fill.background()
        return

    elif tag in ('circle', 'ellipse'):
        cx = float(el.get('cx', 0)); cy = float(el.get('cy', 0))
        if tag == 'circle':
            r = float(el.get('r', 0))
            rx = ry = r
        else:
            rx = float(el.get('rx', 0)); ry = float(el.get('ry', 0))
        x = to_emu(cx - rx, svg_w); y = to_emu(cy - ry, svg_h)
        w = to_emu(rx * 2, svg_w); h = to_emu(ry * 2, svg_h)
        shape = slide.shapes.add_shape(9, x, y, w, h)  # 9 = OVAL
        if sfill:
            shape.fill.solid(); shape.fill.fore_color.rgb = sfill
        else:
            shape.fill.background()
        if sstroke:
            shape.line.color.rgb = sstroke
            shape.line.width = Emu(int(sw * 914400 * 13.333 / svg_w))
        else:
            shape.line.fill.background()
        return

    elif tag == 'line':
        x1 = to_emu(float(el.get('x1', 0)), svg_w)
        y1 = to_emu(float(el.get('y1', 0)), svg_h)
        x2 = to_emu(float(el.get('x2', 0)), svg_w)
        y2 = to_emu(float(el.get('y2', 0)), svg_h)
        w = max(abs(x2 - x1), 1); h = max(abs(y2 - y1), 1)
        x = min(x1, x2); y = min(y1, y2)
        if x2 < x1: x, x2 = x2, x
        if y2 < y1: y, y2 = y2, y
        shape = slide.shapes.add_shape(1, x, y, w, h)
        shape.fill.background()
        if sstroke:
            shape.line.color.rgb = sstroke
            shape.line.width = Emu(int(sw * 914400 * 13.333 / svg_w))
        else:
            shape.line.fill.background()
        return

    elif tag in ('polygon', 'polyline'):
        points_str = el.get('points', '')
        points = []
        for pair in points_str.strip().split():
            if ',' in pair:
                px2, py2 = pair.split(',')
                points.append((float(px2), float(py2)))
        if len(points) < 3:
            return
        # Approximate with bounding box rect for polygon
        xs = [p[0] for p in points]; ys = [p[1] for p in points]
        x = to_emu(min(xs), svg_w); y = to_emu(min(ys), svg_h)
        w = to_emu(max(xs) - min(xs), svg_w); h = to_emu(max(ys) - min(ys), svg_h)
        shape = slide.shapes.add_shape(1, x, y, w, h)
        if sfill:
            shape.fill.solid(); shape.fill.fore_color.rgb = sfill
        else:
            shape.fill.background()
        if sstroke:
            shape.line.color.rgb = sstroke
            shape.line.width = Emu(int(sw * 914400 * 13.333 / svg_w))
        else:
            shape.line.fill.background()
        return

    elif tag == 'path':
        d = el.get('d', '')
        if not d:
            return
        # Try to extract bounding box from path data (simplified)
        # For now, use transform to get bounding box
        # This is a simplified path handler
        try:
            bounds = extract_path_bounds(d)
            if bounds:
                x = to_emu(bounds[0], svg_w); y = to_emu(bounds[1], svg_h)
                w = to_emu(bounds[2] - bounds[0], svg_w); h = to_emu(bounds[3] - bounds[1], svg_h)
                if w > 0 and h > 0:
                    shape = slide.shapes.add_shape(1, x, y, w, h)
                    if sfill:
                        shape.fill.solid(); shape.fill.fore_color.rgb = sfill
                    else:
                        shape.fill.background()
                    if sstroke:
                        shape.line.color.rgb = sstroke
                        shape.line.width = Emu(int(sw * 914400 * 13.333 / svg_w))
                    else:
                        shape.line.fill.background()
        except:
            pass
        return


def extract_path_bounds(d):
    """Extract approximate bounding box from SVG path d string"""
    nums = re.findall(r'[-+]?\d*\.?\d+', d)
    if len(nums) < 4:
        return None
    vals = [float(n) for n in nums]
    return (min(vals[:len(vals)//2*2:2]), min(vals[1:len(vals)//2*2+1:2]),
            max(vals[:len(vals)//2*2:2]), max(vals[1:len(vals)//2*2+1:2]))


def adjust_opacity(color, opacity):
    """Adjust RGB color opacity (simplified - just darken slightly)"""
    if color is None:
        return None
    r = min(255, int(color[0] * opacity + 255 * (1 - opacity)))
    g = min(255, int(color[1] * opacity + 255 * (1 - opacity)))
    b = min(255, int(color[2] * opacity + 255 * (1 - opacity)))
    return RGBColor(r, g, b)


if __name__ == '__main__':
    import sys, os
    svg_file = sys.argv[1] if len(sys.argv) > 1 else 'slides_svg/slide1_cover.svg'
    out_file = sys.argv[2] if len(sys.argv) > 2 else None
    result = svg_to_pptx(svg_file, out_file)
    print(f'Converted: {svg_file} -> {result}')


def create_pptx_with_native_svg(
    svg_files,
    output_path,
    canvas_format='ppt169',
    verbose=False,
    use_compat_mode=True,
    use_native_shapes=True,
    transition='fade',
    transition_duration=2.0,
):
    """Create multi-slide PPTX from multiple SVG files.
    Wraps the single-slide svg_to_pptx logic for batch export."""
    import logging
    logger = logging.getLogger(__name__)

    prs = Presentation()
    prs.slide_width = PPT_W
    prs.slide_height = PPT_H

    for svg_file in svg_files:
        svg_path = str(svg_file) if hasattr(svg_file, 'as_posix') else str(svg_file)
        if verbose:
            logger.info(f'[PPTX Gen] Processing: {svg_path}')

        tree = ET.parse(svg_path)
        root = tree.getroot()

        vb = root.get('viewBox', '0 0 1280 720').split()
        svg_w, svg_h = int(vb[2]), int(vb[3])

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Build defs lookup
        defs = {}
        for defn in root.findall(f'{{{NS}}}defs') or []:
            for el in defn:
                id_ = el.get('id')
                if id_:
                    defs[id_] = el

        # Process graphic elements
        for el in root:
            tag = el.tag.replace(f'{{{NS}}}', '')
            if tag in ('rect', 'circle', 'ellipse', 'line', 'polygon', 'polyline', 'path'):
                process_element(slide, el, tag, defs, svg_w, svg_h)

    output_str = str(output_path) if hasattr(output_path, 'as_posix') else str(output_path)
    os.makedirs(os.path.dirname(output_str), exist_ok=True)
    prs.save(output_str)

    if verbose:
        logger.info(f'[PPTX Gen] Saved: {output_str} ({len(svg_files)} slides)')
    return True
