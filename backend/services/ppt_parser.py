"""
PPT Parser Service - Parse PPTX files to extract design metadata
v0.6 Feature 3: Template configuration - upload PPT, analyze design style
"""
import os
import io
import json
import uuid
import logging
from pathlib import Path
from datetime import datetime

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import PIL.Image
    import PIL.ImageDraw
    import PIL.ImageFont
    HAS_DEPS = True
except ImportError:
    HAS_DEPS = False


# === Color extraction helpers ===

def rgb_to_hex(rgb: tuple) -> str:
    """Convert (r, g, b) to '#RRGGBB'"""
    if rgb is None:
        return '#000000'
    return '#{:02X}{:02X}{:02X}'.format(int(rgb[0]), int(rgb[1]), int(rgb[2]))


def get_slide_background_color(slide) -> str:
    """Extract the dominant background color of a slide."""
    try:
        fill = slide.background.fill
        if fill.type == 1:  # SOLID
            color = fill.fore_color.rgb
            return rgb_to_hex((color.red, color.green, color.blue))
    except Exception:
        pass
    return '#FFFFFF'


def get_text_colors(slide) -> list:
    """Extract distinct text colors from a slide."""
    colors = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    color = run.font.color.rgb
                    if color:
                        hex_c = rgb_to_hex((color.red, color.green, color.blue))
                        if hex_c not in colors:
                            colors.append(hex_c)
                except Exception:
                    pass
    return colors


def get_font_info(slide) -> dict:
    """Extract font information from a slide."""
    fonts = {'title': set(), 'body': set()}
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    fname = run.font.name
                    if fname:
                        size = run.font.size
                        if size and size >= Pt(18):
                            fonts['title'].add(fname)
                        else:
                            fonts['body'].add(fname)
                except Exception:
                    pass
    return {k: list(v) for k, v in fonts.items()}


def get_layout_name(slide) -> str:
    """Get slide layout name."""
    try:
        return slide.slide_layout.name if slide.slide_layout else 'unknown'
    except Exception:
        return 'unknown'


def parse_slide_layout_type(slide) -> str:
    """Infer layout type from shape positions."""
    shapes = list(slide.shapes)
    has_left = any(s.left < Inches(4) for s in shapes if s.left is not None)
    has_right = any(s.left > Inches(4) for s in shapes if s.left is not None)
    if has_left and has_right:
        return '两栏'
    elif shapes:
        # Check if shapes are stacked vertically
        tops = sorted([s.top for s in shapes if s.top is not None])
        if len(tops) >= 2 and tops[-1] - tops[0] > Inches(2):
            return '单栏'
    return '两栏'


def generate_thumbnail(slide, width=400, height=225) -> bytes:
    """Render a slide to PNG bytes using PIL (fallback if no rendering engine)."""
    bg_color = get_slide_background_color(slide)
    img = PIL.Image.new('RGB', (width, height), bg_color or '#FFFFFF')
    draw = PIL.ImageDraw.Draw(img)
    
    # Try to get text content for overlay
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
    
    # Draw placeholder bars to represent content
    bar_y = height // 3
    bar_h = 12
    colors_to_draw = get_text_colors(slide)
    
    if texts:
        draw.rectangle([20, bar_y, width - 20, bar_y + bar_h], fill='#CCCCCC')
        draw.rectangle([20, bar_y + 20, width - 80, bar_y + bar_h + 10], fill='#DDDDDD')
        draw.rectangle([20, bar_y + 40, width - 120, bar_y + bar_h + 8], fill='#DDDDDD')
    
    # Draw a simple header bar at top
    draw.rectangle([0, 0, width, 30], fill='#003371')
    
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    return buf.getvalue()


# === Main Parser Class ===

class PPTParser:
    """
    Parse a PPTX file and extract design metadata.
    """

    def __init__(self, pptx_path: str = None, pptx_bytes: bytes = None):
        if not HAS_DEPS:
            raise RuntimeError("python-pptx or Pillow not installed")
        self.pptx_path = pptx_path
        self.pptx_bytes = pptx_bytes
        self._prs = None

    def _load(self):
        """Load the presentation."""
        if self._prs is not None:
            return
        if self.pptx_path:
            self._prs = Presentation(self.pptx_path)
        elif self.pptx_bytes:
            self._prs = Presentation(io.BytesIO(self.pptx_bytes))
        else:
            raise ValueError("Must provide pptx_path or pptx_bytes")

    def get_slide_count(self) -> int:
        self._load()
        return len(self._prs.slides)

    def get_metadata(self) -> dict:
        """Extract comprehensive metadata from PPTX."""
        self._load()

        slides_meta = []
        all_title_fonts = set()
        all_body_fonts = set()
        all_colors = []

        for i, slide in enumerate(self._prs.slides):
            bg = get_slide_background_color(slide)
            text_colors = get_text_colors(slide)
            fonts = get_font_info(slide)
            layout_name = get_layout_name(slide)
            layout_type = parse_slide_layout_type(slide)

            if bg:
                all_colors.append(bg)
            all_colors.extend(text_colors)
            all_title_fonts.update(fonts['title'])
            all_body_fonts.update(fonts['body'])

            slides_meta.append({
                'index': i,
                'layout_name': layout_name,
                'layout_type': layout_type,
                'bg_color': bg,
                'text_colors': text_colors[:5],  # cap at 5
                'title_fonts': fonts['title'][:3],
                'body_fonts': fonts['body'][:3],
            })

        # Dedupe colors preserving order
        seen = set()
        deduped_colors = []
        for c in all_colors:
            if c not in seen:
                seen.add(c)
                deduped_colors.append(c)

        return {
            'slide_count': len(self._prs.slides),
            'slides': slides_meta,
            'design_meta': {
                'title_fonts': list(all_title_fonts)[:5],
                'body_fonts': list(all_body_fonts)[:5],
                'colors': deduped_colors[:10],
                'layout_type': slides_meta[0]['layout_type'] if slides_meta else '两栏',
            }
        }

    def get_thumbnail(self, slide_index: int, width=400, height=225) -> bytes:
        """Get PNG thumbnail for a specific slide (0-indexed)."""
        self._load()
        if slide_index >= len(self._prs.slides):
            raise IndexError(f"Slide {slide_index} does not exist")
        return generate_thumbnail(self._prs.slides[slide_index], width, height)

    def get_all_thumbnails(self, width=400, height=225) -> dict:
        """Get thumbnails for all slides as {slide_index: bytes}."""
        self._load()
        result = {}
        for i in range(len(self._prs.slides)):
            try:
                result[i] = generate_thumbnail(self._prs.slides[i], width, height)
            except Exception as e:
                logger.warning(f"[PPT Parser] Thumb for slide {i} failed: {e}")
                result[i] = None
        return result